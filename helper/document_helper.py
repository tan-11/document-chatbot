#import PyPDF2
from rank_bm25 import BM25Okapi
import nltk
from nltk.tokenize import word_tokenize
from docx import Document
import re
from nltk.corpus import stopwords
import fitz
import pdfplumber
import io
from PIL import Image
import pytesseract
from sentence_transformers import SentenceTransformer, util
from nltk.tokenize import word_tokenize
from rank_bm25 import BM25Okapi
import torch
from pptx import Presentation
import zipfile
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

nltk.download('punkt')
nltk.download('stopwords')
stop_words = set(stopwords.words('english'))

# Extract text from uploaded PDF
def extract_text_from_pdf(pdf_file):
    """Extracts text and performs OCR on images from a PDF file-like object."""
    all_text = ""

    # Read the PDF file into bytes for consistent reuse
    pdf_bytes = pdf_file.read()
    
    # 1. Extract text using pdfplumber
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            all_text += page_text + "\n"

    # 2. Extract images and apply OCR using PyMuPDF (fitz)
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        for page in doc:
            for img_index, img_info in enumerate(page.get_images(full=True)):
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                try:
                    image = Image.open(io.BytesIO(image_bytes))
                    image_text = pytesseract.image_to_string(image)
                    all_text += image_text + "\n"
                except Exception as e:
                    print(f"Image {img_index} OCR failed: {e}")

    return all_text


def extract_text_from_docx(uploaded_file):
    doc = Document(uploaded_file)
    text = []
    
    # Extract text from paragraphs
    for para in doc.paragraphs:
        if para.text.strip():  # skip empty lines
            text.append(para.text.strip())
    
    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells]
            text.append('\t'.join(row_text))  # Tab-delimited for clarity
    
    try:
        uploaded_file.seek(0)  # Ensure we're reading from the start
        with zipfile.ZipFile(uploaded_file) as docx_zip:
            for file_name in docx_zip.namelist():
                if file_name.startswith("word/media/") and file_name.lower().endswith((".png", ".jpg", ".jpeg")):
                    with docx_zip.open(file_name) as image_file:
                        image = Image.open(io.BytesIO(image_file.read()))
                        image_text = pytesseract.image_to_string(image)
                        text.append(image_text)
    except Exception as e:
        print(f"OCR image extraction failed: {e}")

    return '\n'.join(text)

def extract_text_from_pptx(uploaded_file):
    presentation = Presentation(uploaded_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
            elif shape.shape_type == 13:  # PICTURE type
                try:
                    image_stream = shape.image.blob
                    image = Image.open(io.BytesIO(image_stream))
                    image_text = pytesseract.image_to_string(image)
                    text += image_text + "\n"
                except Exception as e:
                    print(f"OCR failed for PPTX image: {e}")
    return text

def extract_text_from_image(image_file):
    # Set path to Tesseract (Windows)
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

    try:
        image = Image.open(image_file)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"Image OCR failed: {e}")
        return ""


def preprocess_text(text):
    
    # Remove special characters except "."
    text = re.sub(r"[^a-zA-Z0-9.\s]", "", text)
    
    # Tokenize and remove stopwords
    words = word_tokenize(text)
    filtered_words = [word for word in words if word not in stop_words]
    
    # Join the words back into a single string
    preprocessed_text = " ".join(filtered_words)
    return preprocessed_text

def chunk_text(text, words_per_chunk=100):
    words = word_tokenize(text)
    chunks = []
    for i in range(0, len(words), words_per_chunk):
        chunk = words[i:i + words_per_chunk]
        chunks.append(' '.join(chunk))
        
    print(f"Total chunks created: {len(chunks)}")
    return chunks

# Perform BM25 retrieval
def retrieve_relevant_chunks(chunks, query, top_k=2):
    embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    query = query.lower()

    #bm25
    tokenized_chunks = [word_tokenize(chunk.lower()) for chunk in chunks]
    bm25 = BM25Okapi(tokenized_chunks)
    tokenized_query = word_tokenize(query)
    bm25_scores = bm25.get_scores(tokenized_query)

    chunk_embeddings = embedding_model.encode(chunks, convert_to_tensor=True)
    query_embedding = embedding_model.encode(query, convert_to_tensor=True)
    cosine_scores = util.cos_sim(query_embedding, chunk_embeddings)[0]

    bm25_tensor = torch.tensor(bm25_scores)
    bm25_tensor = (bm25_tensor - bm25_tensor.min()) / (bm25_tensor.max() - bm25_tensor.min() + 1e-8)
    combined_scores = 0.5 * cosine_scores + 0.5 * bm25_tensor


    top_chunks = torch.topk(combined_scores, min(top_k, len(chunks))).indices
    extended_top_indices = set()
    for i in top_chunks:
        extended_top_indices.add(i.item())
        if i - 1 >= 0:
            extended_top_indices.add(i - 1)
        if i + 1 < len(chunks):
            extended_top_indices.add(i + 1)
    print(f"Extended top indices: {extended_top_indices}")
    return [chunks[i] for i in extended_top_indices]
